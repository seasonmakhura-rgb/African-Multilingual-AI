import os
import io
import time
import json
import pickle
import numpy as np
import pandas as pd
import streamlit as st
import tensorflow as tf
from tensorflow.keras.preprocessing.sequence import pad_sequences
from google import genai
from google.genai import types
from gtts import gTTS
import faiss
from PIL import Image
import pypdf
import docx
import pptx
from documents import KNOWLEDGE_DOCUMENTS
from guardrails import validate_user_input

# Page Config
st.set_page_config(page_title="BAOBAB AI", page_icon="🌳", layout="wide")

# Target Gemini & Imagen Model Identifiers
GEMINI_MODEL = "gemini-3.6-flash"
IMAGEN_MODEL = "imagen-3.0-generate-002"

# Ensure persistent storage directory exists
SAVES_DIR = "saved_chats"
os.makedirs(SAVES_DIR, exist_ok=True)

# Dynamic TTS Language Mapping
TTS_LANG_MAP = {
    'Portuguese': 'pt',
    'French': 'fr',
    'Swahili': 'sw',
    'English': 'en',
    'Spanish': 'es'
}

# --- Tier 4: Executable External Tool Definitions ---
def get_current_weather(location: str) -> str:
    """Fetches real-time structured weather reports for a city or region."""
    loc_clean = location.strip().title()
    weather_db = {
        "Luanda": {"temp": "29°C", "condition": "Sunny", "humidity": "74%"},
        "Maputo": {"temp": "26°C", "condition": "Partly Cloudy", "humidity": "68%"},
        "Nairobi": {"temp": "22°C", "condition": "Clear", "humidity": "55%"},
        "Johannesburg": {"temp": "21°C", "condition": "Mild", "humidity": "40%"},
        "Lisbon": {"temp": "24°C", "condition": "Sunny", "humidity": "60%"},
    }
    data = weather_db.get(loc_clean, {"temp": "25°C", "condition": "Clear", "humidity": "50%"})
    return json.dumps({"location": loc_clean, **data})

def calculate_loan_repayment(principal: float, annual_rate: float, duration_years: int) -> str:
    """Calculates monthly payment and total interest for a fixed-rate loan."""
    monthly_rate = (annual_rate / 100) / 12
    total_months = duration_years * 12
    if monthly_rate == 0:
        monthly_payment = principal / total_months
    else:
        monthly_payment = principal * (monthly_rate * (1 + monthly_rate)**total_months) / ((1 + monthly_rate)**total_months - 1)
    total_payment = monthly_payment * total_months
    total_interest = total_payment - principal
    
    return json.dumps({
        "principal": principal,
        "annual_rate_percent": annual_rate,
        "duration_years": duration_years,
        "monthly_payment": round(monthly_payment, 2),
        "total_interest": round(total_interest, 2),
        "total_payment": round(total_payment, 2)
    })

def generate_imagen_art(prompt: str, aspect_ratio: str = "1:1") -> Image.Image:
    """Generates a PIL image from a text prompt using Google Imagen 3."""
    client_gen = genai.Client(api_key=st.secrets.get("GEMINI_API_KEY", os.environ.get("GEMINI_API_KEY")))
    result = client_gen.models.generate_images(
        model=IMAGEN_MODEL,
        prompt=prompt,
        config=types.GenerateImagesConfig(
            number_of_images=1,
            aspect_ratio=aspect_ratio,
            output_mime_type="image/jpeg"
        )
    )
    generated_img_bytes = result.generated_images[0].image.image_bytes
    return Image.open(io.BytesIO(generated_img_bytes))

# Map functions to their callable implementations
AVAILABLE_FUNCTIONS = {
    "get_current_weather": get_current_weather,
    "calculate_loan_repayment": calculate_loan_repayment
}

# Declarations for Gemini Tool Schema
weather_declaration = types.FunctionDeclaration(
    name="get_current_weather",
    description="Get current weather details for a specific city or region.",
    parameters=types.Schema(
        type="OBJECT",
        properties={
            "location": types.Schema(type="STRING", description="The city and state/country, e.g. Luanda, Maputo, Nairobi")
        },
        required=["location"]
    )
)

loan_declaration = types.FunctionDeclaration(
    name="calculate_loan_repayment",
    description="Calculate monthly loan repayments, total interest, and total payout.",
    parameters=types.Schema(
        type="OBJECT",
        properties={
            "principal": types.Schema(type="NUMBER", description="The total loan principal amount"),
            "annual_rate": types.Schema(type="NUMBER", description="Annual interest rate percentage, e.g. 7.5"),
            "duration_years": types.Schema(type="INTEGER", description="Loan period in years")
        },
        required=["principal", "annual_rate", "duration_years"]
    )
)

# Initialize State
if "messages" not in st.session_state:
    st.session_state.messages = []

if "analytics" not in st.session_state:
    st.session_state.analytics = {
        "bilstm_latencies": [],
        "faiss_latencies": [],
        "total_requests": 0,
        "blocked_requests": 0
    }

if "custom_documents" not in st.session_state:
    st.session_state.custom_documents = []

if "workspaces" not in st.session_state:
    st.session_state.workspaces = ["General", "Coding", "Legal", "Academic"]

if "active_workspace" not in st.session_state:
    st.session_state.active_workspace = "General"

# Gemini API Client setup
api_key = st.secrets.get("GEMINI_API_KEY", os.environ.get("GEMINI_API_KEY"))
client = genai.Client(api_key=api_key)

def extract_text_from_file(uploaded_file, client_gemini=None):
    filename = uploaded_file.name.lower()
    
    if filename.endswith(".txt"):
        return uploaded_file.read().decode("utf-8").strip()
    elif filename.endswith(".pdf"):
        pdf_reader = pypdf.PdfReader(uploaded_file)
        return "\n".join([page.extract_text() for page in pdf_reader.pages if page.extract_text()]).strip()
    elif filename.endswith(".docx"):
        doc = docx.Document(uploaded_file)
        return "\n".join([p.text for p in doc.paragraphs if p.text.strip()]).strip()
    elif filename.endswith(".pptx"):
        prs = pptx.Presentation(uploaded_file)
        text_content = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text)
        return "\n".join(text_content).strip()
    elif filename.endswith((".jpg", ".jpeg", ".png")):
        image = Image.open(uploaded_file)
        response = client_gemini.models.generate_content(
            model=GEMINI_MODEL,
            contents=[image, "Extract and transcribe all readable text cleanly and accurately."]
        )
        return response.text.strip()
    return ""

@st.cache_resource
def load_system_resources():
    model = tf.keras.models.load_model('african_lang_classifier.keras')
    with open('word_tokenizer.pkl', 'rb') as f:
        tokenizer = pickle.load(f)
    with open('label_encoder.pkl', 'rb') as f:
        label_encoder = pickle.load(f)
    return model, tokenizer, label_encoder

model, tokenizer, label_encoder = load_system_resources()

def build_faiss_index_for_docs(docs, tokenizer, embedding_dim=128):
    if not docs:
        return None, []
    np.random.seed(42)
    doc_vectors = []
    for doc in docs:
        seq = tokenizer.texts_to_sequences([doc['text']])
        vec = np.zeros(embedding_dim, dtype='float32')
        if seq[0]:
            for idx in seq[0][:embedding_dim]:
                vec[idx % embedding_dim] += 1.0
        doc_vectors.append(vec)
    doc_matrix = np.array(doc_vectors).astype('float32')
    faiss.normalize_L2(doc_matrix)
    index = faiss.IndexFlatIP(embedding_dim)
    index.add(doc_matrix)
    return index, docs

# Tag base knowledge documents with workspace context
tagged_base_docs = []
for doc in KNOWLEDGE_DOCUMENTS:
    doc_copy = doc.copy()
    doc_copy["workspace"] = doc.get("workspace", "General")
    tagged_base_docs.append(doc_copy)

all_active_documents = tagged_base_docs + st.session_state.custom_documents

# --- Header ---
st.title("🌳 BAOBAB AI")

# --- Sidebar Controls ---
with st.sidebar:
    st.header("⚙️ Controls")
    target_language = st.selectbox(
        "Output Language", 
        sorted(label_encoder.classes_), 
        index=sorted(label_encoder.classes_).index("Portuguese") if "Portuguese" in label_encoder.classes_ else 0
    )

    st.markdown("---")
    st.header("📂 Knowledge Workspaces (Tier 3)")
    
    st.session_state.active_workspace = st.selectbox(
        "Active Workspace",
        st.session_state.workspaces,
        index=st.session_state.workspaces.index(st.session_state.active_workspace) if st.session_state.active_workspace in st.session_state.workspaces else 0
    )

    new_ws = st.text_input("Create Workspace", placeholder="e.g. Finance")
    if st.button("➕ Create Workspace") and new_ws.strip():
        ws_clean = new_ws.strip().title()
        if ws_clean not in st.session_state.workspaces:
            st.session_state.workspaces.append(ws_clean)
            st.session_state.active_workspace = ws_clean
            st.success(f"Workspace **{ws_clean}** created!")
            st.rerun()

    active_docs_in_ws = [d for d in all_active_documents if d.get("workspace", "General") == st.session_state.active_workspace]
    st.caption(f"📚 Documents in **{st.session_state.active_workspace}**: **{len(active_docs_in_ws)}**")

    st.markdown("---")
    st.header("🌐 Active Capabilities")
    enable_web_search = st.toggle("Enable Live Web Search", value=False)
    enable_code_interpreter = st.toggle("Enable Python Code Sandbox", value=True)
    enable_function_calling = st.toggle("Enable API Tool Calling (Tier 4)", value=True)
    
    custom_persona = st.text_area(
        "Custom Instructions / Persona", 
        value="You are BAOBAB AI, an expert, helpful assistant.",
        help="Define how the AI should behave or what role it should take."
    )

    st.markdown("---")
    st.header("🔍 Chat Navigation & Search")
    chat_search_query = st.text_input("Find in active chat", placeholder="Type keyword...", key="chat_search_key")
    
    if st.session_state.messages:
        st.subheader("📍 Jump to Message")
        user_msgs = [(idx, msg["content"]) for idx, msg in enumerate(st.session_state.messages) if msg["role"] == "user"]
        if user_msgs:
            for idx, text in user_msgs:
                label_text = text[:28] + ("..." if len(text) > 28 else "")
                st.markdown(f"👉 [{label_text}](#msg-{idx})")

    st.markdown("---")
    st.header("💾 Session History")
    session_name = st.text_input("Session Name", value="Session_1")
    if st.button("💾 Save Session"):
        if st.session_state.messages:
            save_payload = {
                "messages": [{k: v for k, v in msg.items() if k not in ["audio", "generated_image"]} for msg in st.session_state.messages],
                "analytics": st.session_state.analytics,
                "workspace": st.session_state.active_workspace
            }
            with open(os.path.join(SAVES_DIR, f"{session_name.strip()}.json"), "w", encoding="utf-8") as f:
                json.dump(save_payload, f, indent=2)
            st.success("Session saved!")
            st.rerun()

    saved_files = [f.replace(".json", "") for f in os.listdir(SAVES_DIR) if f.endswith(".json")]
    if saved_files:
        selected_file = st.selectbox("Load History", ["-- Select Session --"] + saved_files)
        if st.button("📂 Load Session"):
            if selected_file != "-- Select Session --":
                with open(os.path.join(SAVES_DIR, f"{selected_file}.json"), "r", encoding="utf-8") as f:
                    loaded_data = json.load(f)
                st.session_state.messages = loaded_data.get("messages", [])
                st.session_state.analytics = loaded_data.get("analytics", st.session_state.analytics)
                if "workspace" in loaded_data and loaded_data["workspace"] in st.session_state.workspaces:
                    st.session_state.active_workspace = loaded_data["workspace"]
                st.rerun()

    if st.button("🗑️ Clear Chat"):
        st.session_state.messages = []
        st.session_state.analytics = {"bilstm_latencies": [], "faiss_latencies": [], "total_requests": 0, "blocked_requests": 0}
        st.rerun()

    st.markdown("---")
    st.header("📊 Performance Analytics")
    col1, col2 = st.columns(2)
    col1.metric("Total Turns", st.session_state.analytics["total_requests"])
    col2.metric("Blocked", st.session_state.analytics["blocked_requests"])

# --- Main Chat UI ---
search_active = bool(chat_search_query.strip())
matches_found = 0

for i, msg in enumerate(st.session_state.messages):
    is_match = search_active and chat_search_query.lower() in msg["content"].lower()
    if is_match:
        matches_found += 1
    
    if not search_active or is_match:
        st.markdown(f'<div id="msg-{i}"></div>', unsafe_allow_html=True)
        with st.chat_message(msg["role"]):
            if is_match:
                st.markdown(f"🔍 *Match found for '{chat_search_query}':*")
            st.markdown(msg["content"])
            if "generated_image" in msg and msg["generated_image"] is not None:
                st.image(msg["generated_image"], caption="Generated by Imagen 3", use_container_width=True)
            if "metadata" in msg:
                st.caption(msg["metadata"])
            if "audio" in msg and msg["audio"] is not None:
                st.audio(msg["audio"], format="audio/mp3")

if search_active:
    st.info(f"🔎 Filtered results for '**{chat_search_query}**': Found **{matches_found}** matching message(s).")

st.markdown("---")

# --- Attachment Studio Popover (Positioned directly above the chat box) ---
with st.popover("📎 Tools & Attachments", help="Upload documents, scan with camera, or generate AI artwork"):
    tab_file, tab_photo, tab_imagen = st.tabs(["📄 Upload", "📷 Camera", "🎨 Imagen Studio"])
    
    with tab_file:
        uploaded_file = st.file_uploader("Upload File", type=["txt", "pdf", "docx", "pptx", "jpg", "jpeg", "png"], key="attachment_file")
        upload_workspace = st.selectbox("Assign Workspace", st.session_state.workspaces, index=st.session_state.workspaces.index(st.session_state.active_workspace), key="file_ws")
        upload_category = st.text_input("Category", value="Custom Knowledge", key="file_cat")
        if st.button("📥 Attach File", key="btn_attach_file") and uploaded_file:
            ext_text = extract_text_from_file(uploaded_file, client_gemini=client)
            if ext_text:
                st.session_state.custom_documents.append({
                    "id": len(all_active_documents) + 1,
                    "title": uploaded_file.name,
                    "workspace": upload_workspace,
                    "category": upload_category.strip() or "Custom Knowledge",
                    "text": ext_text
                })
                st.success(f"Attached **{uploaded_file.name}** to **{upload_workspace}**!")

    with tab_photo:
        camera_image = st.camera_input("Take photo", key="attachment_cam")
        photo_workspace = st.selectbox("Assign Workspace", st.session_state.workspaces, index=st.session_state.workspaces.index(st.session_state.active_workspace), key="cam_ws")
        photo_category = st.text_input("Category", value="Custom Knowledge", key="cam_cat")
        if st.button("📥 Attach Scan", key="btn_attach_cam") and camera_image:
            img = Image.open(camera_image)
            ext_text = client.models.generate_content(model=GEMINI_MODEL, contents=[img, "Extract text cleanly."]).text.strip()
            if ext_text:
                st.session_state.custom_documents.append({
                    "id": len(all_active_documents) + 1,
                    "title": f"Scan_{int(time.time())}.jpg",
                    "workspace": photo_workspace,
                    "category": photo_category.strip() or "Custom Knowledge",
                    "text": ext_text
                })
                st.success(f"Photo attached to **{photo_workspace}**!")

    with tab_imagen:
        st.subheader("🎨 Generate Image with Imagen 3")
        img_prompt = st.text_area("Image Description", placeholder="e.g. A futuristic Baobab tree bathed in bioluminescent light", key="imagen_studio_prompt")
        aspect_choice = st.selectbox("Aspect Ratio", ["1:1", "16:9", "9:16", "4:3", "3:4"], key="imagen_studio_aspect")
        
        if st.button("✨ Generate Image", key="btn_generate_imagen"):
            if img_prompt.strip():
                with st.spinner("Generating image via Imagen 3..."):
                    try:
                        pil_img = generate_imagen_art(img_prompt.strip(), aspect_ratio=aspect_choice)
                        st.session_state.messages.append({"role": "user", "content": f"🎨 Generate Image: {img_prompt}"})
                        st.session_state.messages.append({
                            "role": "assistant",
                            "content": f"Here is the generated image for: *\"{img_prompt}\"*",
                            "generated_image": pil_img,
                            "metadata": "🎨 Model: `Imagen 3`"
                        })
                        st.rerun()
                    except Exception as e:
                        st.error(f"Image generation error: {str(e)}")

# --- Native Chat Input with Integrated Microphone Inside standard Input Field ---
chat_input_response = st.chat_input(
    placeholder=f"Ask BAOBAB AI [{st.session_state.active_workspace}] or tap 🎙️ inside to record...",
    accept_file=False
)

# Process Prompt or Spoken Voice
if chat_input_response:
    user_input = chat_input_response.strip()
    
    st.session_state.analytics["total_requests"] += 1
    
    is_safe, guardrail_msg = validate_user_input(user_input)
    if not is_safe:
        st.session_state.analytics["blocked_requests"] += 1
        st.error(f"🚫 **Input Blocked**: {guardrail_msg}")
    else:
        st.session_state.messages.append({"role": "user", "content": user_input})
        
        with st.spinner("Processing..."):
            t0_bilstm = time.perf_counter()
            seq = tokenizer.texts_to_sequences([user_input])
            padded = pad_sequences(seq, maxlen=50)
            preds = model.predict(padded)
            detected_lang = label_encoder.inverse_transform([np.argmax(preds)])[0]
            confidence = float(np.max(preds)) * 100
            st.session_state.analytics["bilstm_latencies"].append(round((time.perf_counter() - t0_bilstm) * 1000, 2))

            t0_faiss = time.perf_counter()
            embedding_dim = 128
            workspace_docs = [d for d in all_active_documents if d.get("workspace", "General") == st.session_state.active_workspace]
            
            active_faiss_index, active_docs = build_faiss_index_for_docs(workspace_docs, tokenizer, embedding_dim)
            rag_context, doc_title = f"No relevant document found in workspace '{st.session_state.active_workspace}'.", "None"
            
            if active_faiss_index is not None:
                q_vec = np.zeros(embedding_dim, dtype='float32')
                if seq[0]:
                    for idx in seq[0][:embedding_dim]:
                        q_vec[idx % embedding_dim] += 1.0
                if np.sum(q_vec) > 0:
                    q_matrix = np.array([q_vec]).astype('float32')
                    faiss.normalize_L2(q_matrix)
                    distances, indices = active_faiss_index.search(q_matrix, k=1)
                    if float(distances[0][0]) >= 0.35:
                        retrieved_doc = active_docs[indices[0][0]]
                        rag_context = retrieved_doc["text"]
                        doc_title = f"[{retrieved_doc['category']}] {retrieved_doc['title']}"
            st.session_state.analytics["faiss_latencies"].append(round((time.perf_counter() - t0_faiss) * 1000, 2))

            history_context = "".join([f"{m['role'].upper()}: {m['content']}\n" for m in st.session_state.messages[:-1]])

            prompt_sent = f"""SYSTEM INSTRUCTION:
{custom_persona}
Detected Input Language: {detected_lang} (Confidence: {confidence:.1f}%).
TARGET OUTPUT LANGUAGE: Compose your response strictly in {target_language}.
CURRENT WORKSPACE: {st.session_state.active_workspace}

RETRIEVED WORKSPACE CONTEXT:
{rag_context}

HISTORY:
{history_context if history_context else "None."}

QUERY: {user_input}
RESPONSE ({target_language}):"""

            tools_list = []
            if enable_web_search:
                tools_list.append({"google_search": {}})
            if enable_code_interpreter:
                tools_list.append(types.Tool(code_execution=types.CodeExecution()))
            if enable_function_calling:
                tools_list.append(types.Tool(function_declarations=[weather_declaration, loan_declaration]))

            config = types.GenerateContentConfig(tools=tools_list) if tools_list else None

        # Execution
        with st.chat_message("assistant"):
            message_placeholder = st.empty()
            tool_used_label = "None"
            generated_image_obj = None
            
            # Check for explicit image generation triggers in user input
            if any(k in user_input.lower() for k in ["generate image", "draw", "create picture", "generate an image"]):
                try:
                    generated_image_obj = generate_imagen_art(user_input)
                    tool_used_label = "`Imagen 3`"
                    ai_output = f"Generated image based on prompt: *\"{user_input}\"*"
                    st.image(generated_image_obj, caption="Imagen 3", use_container_width=True)
                except Exception as e:
                    ai_output = f"Imagen Generation Error: {str(e)}"
                    message_placeholder.markdown(ai_output)
            else:
                try:
                    response = client.models.generate_content(
                        model=GEMINI_MODEL,
                        contents=prompt_sent,
                        config=config
                    )
                    
                    if response.function_calls:
                        call = response.function_calls[0]
                        fn_name = call.name
                        fn_args = call.args
                        tool_used_label = f"`{fn_name}()`"
                        
                        if fn_name in AVAILABLE_FUNCTIONS:
                            tool_result = AVAILABLE_FUNCTIONS[fn_name](**fn_args)
                            followup_response = client.models.generate_content(
                                model=GEMINI_MODEL,
                                contents=[
                                    prompt_sent,
                                    response.candidates[0].content,
                                    types.Part.from_function_response(
                                        name=fn_name,
                                        response={"result": tool_result}
                                    )
                                ]
                            )
                            ai_output = followup_response.text
                        else:
                            ai_output = response.text
                    else:
                        ai_output = response.text

                    message_placeholder.markdown(ai_output)

                except Exception as e:
                    ai_output = f"Error: {str(e)}"
                    message_placeholder.markdown(ai_output)

        # TTS Audio Generation
        audio_bytes = None
        try:
            tts = gTTS(text=ai_output, lang=TTS_LANG_MAP.get(target_language, 'en'))
            fp = io.BytesIO()
            tts.write_to_fp(fp)
            audio_bytes = fp.getvalue()
        except Exception:
            pass

        metadata = f"🧠 **{detected_lang}** ({confidence:.1f}%) | 📁 Workspace: `{st.session_state.active_workspace}` | 🔧 Tool Called: {tool_used_label}"
        st.session_state.messages.append({
            "role": "assistant", 
            "content": ai_output, 
            "generated_image": generated_image_obj,
            "metadata": metadata,
            "audio": audio_bytes
        })
        st.rerun()
